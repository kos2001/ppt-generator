#!/usr/bin/env python3
"""Map the addressable shapes of an existing .pptx so edits can target them.

Editing an existing deck requires knowing WHAT to point at. This prints, per
slide, every shape with its index (its position in slide.shapes — the address
edit_pptx.py uses), name, type, a text preview, and table/image dimensions.

Usage:
    python inspect_pptx.py deck.pptx           # human-readable map
    python inspect_pptx.py deck.pptx --json     # machine-readable

Reads ordinary, unprotected .pptx files only. A DRM-locked file must first be
exported to a plaintext copy through its DRM client by an authorized user.
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def _type_name(shape):
    try:
        return shape.shape_type.name if shape.shape_type is not None else "UNKNOWN"
    except Exception:  # noqa: BLE001 - some placeholders raise on shape_type
        return "PLACEHOLDER"


def describe_shape(shape, idx):
    info = {
        "shape": idx,
        "name": shape.name,
        "type": _type_name(shape),
    }
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            info["text"] = text
    if shape.has_table:
        tbl = shape.table
        info["table"] = {"rows": len(tbl.rows), "cols": len(tbl.columns)}
    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "series": len(chart.series),
            "categories": len(list(chart.plots[0].categories)) if chart.plots else 0,
        }
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["image"] = True
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["group"] = [describe_shape(s, i) for i, s in enumerate(shape.shapes)]
    return info


def inspect(path):
    prs = Presentation(path)
    slides = []
    for sidx, slide in enumerate(prs.slides, start=1):
        shapes = [describe_shape(sh, i) for i, sh in enumerate(slide.shapes)]
        slides.append({"slide": sidx, "shapes": shapes})
    return slides


def to_text(slides):
    out = []
    for s in slides:
        out.append("Slide %d" % s["slide"])
        for sh in s["shapes"]:
            line = "  [%d] %-9s name=%r" % (sh["shape"], sh["type"], sh["name"])
            if "table" in sh:
                line += " table=%dx%d" % (sh["table"]["rows"], sh["table"]["cols"])
            if "chart" in sh:
                line += " chart=%s(%dser x %dcat)" % (
                    sh["chart"]["chart_type"].split(" ")[0],
                    sh["chart"]["series"], sh["chart"]["categories"])
            if sh.get("image"):
                line += " image"
            out.append(line)
            if "text" in sh:
                preview = sh["text"].replace("\n", " / ")
                if len(preview) > 70:
                    preview = preview[:67] + "..."
                out.append("        text: %s" % preview)
        out.append("")
    return "\n".join(out)


def main(argv=None):
    ap = argparse.ArgumentParser(description="Map addressable shapes in a .pptx.")
    ap.add_argument("path", help="path to .pptx file")
    ap.add_argument("--json", action="store_true", help="output JSON")
    args = ap.parse_args(argv)

    try:
        slides = inspect(args.path)
    except Exception as e:  # noqa: BLE001
        print("Failed to read '%s': %s" % (args.path, e), file=sys.stderr)
        return 1

    if args.json:
        print(json.dumps(slides, ensure_ascii=False, indent=2))
    else:
        print(to_text(slides))
    return 0


if __name__ == "__main__":
    sys.exit(main())

#!/usr/bin/env python
"""Extract text, tables, and speaker notes from a (DRM-free) PPTX file.

Usage:
    python scripts/extract_pptx.py <path-to.pptx> [--json]

Note: This reads ordinary, unprotected .pptx files only. A DRM-locked file
must first be exported to a plaintext copy through its DRM client by an
authorized user — this script does not and will not bypass DRM.
"""
import argparse
import json
import os
import sys

from pptx import Presentation

# Ensure UTF-8 output on consoles with a legacy default codec (e.g. cp949).
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def iter_items(shape, slide_idx, img_dir, counter):
    """Yield text/table/image items from a shape, recursing into groups.

    Embedded images are written to *img_dir* (if given) and reported by path.
    """
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            yield from iter_items(sub, slide_idx, img_dir, counter)
        return
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            yield {"kind": "text", "text": text}
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        yield {"kind": "table", "rows": rows}
    if shape.has_chart:
        chart = shape.chart
        try:
            categories = [str(c) for c in chart.plots[0].categories]
        except (IndexError, ValueError):
            categories = []
        series = [{"name": s.name, "values": list(s.values)} for s in chart.series]
        yield {
            "kind": "chart",
            "chart_type": str(chart.chart_type),
            "categories": categories,
            "series": series,
        }
    if shape.shape_type == 13:  # PICTURE
        try:
            image = shape.image
        except Exception:  # noqa: BLE001 - shape without an extractable blob
            return
        counter[0] += 1
        name = f"slide{slide_idx}_img{counter[0]}.{image.ext}"
        item = {"kind": "image", "filename": name, "content_type": image.content_type}
        if img_dir:
            out_path = os.path.join(img_dir, name)
            with open(out_path, "wb") as f:
                f.write(image.blob)
            item["path"] = out_path
        yield item


def extract(path, img_dir=None):
    if img_dir:
        os.makedirs(img_dir, exist_ok=True)
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        items = []
        counter = [0]
        for shape in slide.shapes:
            items.extend(iter_items(shape, idx, img_dir, counter))
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"slide": idx, "items": items, "notes": notes})
    return slides


def to_markdown(slides):
    out = []
    for s in slides:
        out.append(f"## Slide {s['slide']}")
        for item in s["items"]:
            if item["kind"] == "text":
                out.append(item["text"])
            elif item["kind"] == "table":
                rows = item["rows"]
                if not rows:
                    continue
                out.append("| " + " | ".join(rows[0]) + " |")
                out.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for r in rows[1:]:
                    out.append("| " + " | ".join(r) + " |")
            elif item["kind"] == "image":
                ref = item.get("path", item["filename"])
                out.append(f"![{item['filename']}]({ref})")
            elif item["kind"] == "chart":
                out.append(f"**Chart** ({item['chart_type']})")
                cats = item["categories"]
                header = ["series"] + (cats if cats else [])
                out.append("| " + " | ".join(header) + " |")
                out.append("| " + " | ".join("---" for _ in header) + " |")
                for ser in item["series"]:
                    vals = ["" if v is None else str(v) for v in ser["values"]]
                    out.append("| " + " | ".join([str(ser["name"])] + vals) + " |")
        if s["notes"]:
            out.append(f"\n> **Notes:** {s['notes']}")
        out.append("")
    return "\n".join(out)


def main():
    ap = argparse.ArgumentParser(description="Extract content from a DRM-free PPTX.")
    ap.add_argument("path", help="path to .pptx file")
    ap.add_argument("--json", action="store_true", help="output JSON instead of markdown")
    ap.add_argument(
        "--images", metavar="DIR",
        help="extract embedded images into DIR (omit to skip image files)",
    )
    args = ap.parse_args()

    try:
        slides = extract(args.path, img_dir=args.images)
    except Exception as e:  # noqa: BLE001
        print(f"Failed to read '{args.path}': {e}", file=sys.stderr)
        print(
            "If this is DRM-protected, export a plaintext copy via your DRM "
            "client first (this tool cannot read protected files).",
            file=sys.stderr,
        )
        sys.exit(1)

    if args.json:
        print(json.dumps(slides, ensure_ascii=False, indent=2))
    else:
        print(to_markdown(slides))


if __name__ == "__main__":
    main()

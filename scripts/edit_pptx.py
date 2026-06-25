#!/usr/bin/env python3
"""Apply a JSON edit spec to an existing .pptx, preserving its design.

This is the counterpart to build_pptx.py: instead of rendering a new deck from
scratch, it opens an existing presentation and applies surgical edits — keeping
every untouched shape, master, image, and bit of formatting exactly as it was.

The edit spec lists operations addressed at shapes by (slide, shape) index as
reported by inspect_pptx.py. Supported operations:

    replace_text   find/replace across the whole deck or given slides,
                   preserving each run's character formatting (font, color,
                   bold) — the naive `shape.text = ...` collapses runs and
                   loses all of it, so this walks runs instead.
    set_text       replace all text of one shape (keeps the first run's font)
    set_table_cell set one table cell's text by row/col
    replace_image  swap a picture's image bytes, keeping its position/size
    add_image      add a new picture to a slide at an inch-based position
    delete_shape   remove a shape

Usage:
    python edit_pptx.py edits.json                 # source/output from spec
    python edit_pptx.py edits.json -i in.pptx -o out.pptx

Reads ordinary, unprotected .pptx files only. A DRM-locked file must first be
exported to a plaintext copy through its DRM client by an authorized user; this
tool does not bypass DRM.
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


class EditError(Exception):
    pass


# --------------------------------------------------------------------------- #
# Shape addressing
# --------------------------------------------------------------------------- #
def _slide(prs, n):
    slides = list(prs.slides)
    if not (1 <= n <= len(slides)):
        raise EditError("slide %d out of range (deck has %d)" % (n, len(slides)))
    return slides[n - 1]


def _shape(prs, slide_n, shape_n):
    slide = _slide(prs, slide_n)
    shapes = list(slide.shapes)
    if not (0 <= shape_n < len(shapes)):
        raise EditError("slide %d shape %d out of range (slide has %d shapes)"
                        % (slide_n, shape_n, len(shapes)))
    return slide, shapes[shape_n]


# --------------------------------------------------------------------------- #
# Run-preserving text replacement
# --------------------------------------------------------------------------- #
def _replace_in_paragraph(para, find, replace):
    """Replace every occurrence of *find* within a paragraph, even when it is
    split across runs, keeping the formatting of the run where each match
    starts. Returns the number of replacements made."""
    runs = para.runs
    if not runs or not find:
        return 0
    full = "".join(r.text for r in runs)
    if find not in full:
        return 0

    # Map each character position in `full` to its originating run index.
    owner = []
    for i, r in enumerate(runs):
        owner.extend([i] * len(r.text))

    count = 0
    start = full.find(find)
    while start != -1:
        end = start + len(find)
        first = owner[start]
        # For each run, drop the characters overlapping [start, end); the
        # replacement text is inserted into the run where the match begins, so
        # it inherits that run's font/color/bold.
        result = {}
        cursor = 0
        for i, r in enumerate(runs):
            seg = r.text
            seg_start, seg_end = cursor, cursor + len(seg)
            cursor = seg_end
            keep_before = seg[: max(0, min(seg_end, start) - seg_start)] if seg_start < start else ""
            keep_after = seg[max(0, end - seg_start):] if seg_end > end else ""
            result[i] = keep_before + (replace if i == first else "") + keep_after
        for i, r in enumerate(runs):
            r.text = result[i]
        count += 1
        # Recompute and look for the next occurrence.
        full = "".join(r.text for r in runs)
        owner = []
        for i, r in enumerate(runs):
            owner.extend([i] * len(r.text))
        start = full.find(find, 0)
    return count


def _iter_text_frames(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from _iter_text_frames(sh.shapes)
            continue
        if sh.has_text_frame:
            yield sh.text_frame
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def op_replace_text(prs, op):
    find = op["find"]
    replace = op.get("replace", "")
    slides = op.get("slides")  # list of 1-based slide numbers, or None for all
    count = 0
    for sidx, slide in enumerate(prs.slides, start=1):
        if slides and sidx not in slides:
            continue
        for tf in _iter_text_frames(slide.shapes):
            for para in tf.paragraphs:
                count += _replace_in_paragraph(para, find, replace)
    return "replace_text %r -> %r: %d replacement(s)" % (find, replace, count)


def op_set_text(prs, op):
    _, shape = _shape(prs, op["slide"], op["shape"])
    if not shape.has_text_frame:
        raise EditError("slide %d shape %d has no text frame" % (op["slide"], op["shape"]))
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = op["text"]
        for r in para.runs[1:]:
            r.text = ""
    else:
        para.text = op["text"]
    # Drop any extra paragraphs so only the new text remains.
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    return "set_text slide %d shape %d" % (op["slide"], op["shape"])


def op_set_table_cell(prs, op):
    _, shape = _shape(prs, op["slide"], op["shape"])
    if not shape.has_table:
        raise EditError("slide %d shape %d is not a table" % (op["slide"], op["shape"]))
    cell = shape.table.cell(op["row"], op["col"])
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = op["text"]
        for r in para.runs[1:]:
            r.text = ""
    else:
        cell.text = op["text"]
    return "set_table_cell slide %d shape %d [%d,%d]" % (
        op["slide"], op["shape"], op["row"], op["col"])


def _image_px(path):
    """Return (width, height) in pixels, or None if it can't be read."""
    try:
        from PIL import Image
        with Image.open(path) as im:
            return im.size
    except Exception:  # noqa: BLE001
        return None


def _apply_crop(pic, crop_l, crop_t, crop_r, crop_b):
    pic.crop_left = crop_l
    pic.crop_top = crop_t
    pic.crop_right = crop_r
    pic.crop_bottom = crop_b


def _fit_into_box(pic, iw, ih, box_l, box_t, box_w, box_h, mode):
    """Resize/position an already-added picture to a fixed box.

    contain: scale to fit inside the box, centered, no crop (letterboxed).
    cover:   fill the box exactly, cropping the overflow, no distortion.
    stretch: fill the box exactly, distorting to match (no crop).
    """
    if mode == "stretch" or not iw or not ih:
        pic.left, pic.top, pic.width, pic.height = box_l, box_t, box_w, box_h
        return
    box_ar = box_w / box_h
    img_ar = iw / ih
    if mode == "cover":
        # Picture frame == box; crop the axis that overflows so visible area
        # matches the box aspect with no distortion.
        pic.left, pic.top, pic.width, pic.height = box_l, box_t, box_w, box_h
        if img_ar > box_ar:                      # image too wide -> crop sides
            frac = (1 - box_ar / img_ar) / 2
            _apply_crop(pic, frac, 0, frac, 0)
        else:                                     # image too tall -> crop top/bottom
            frac = (1 - img_ar / box_ar) / 2
            _apply_crop(pic, 0, frac, 0, frac)
    else:  # contain (default)
        scale = min(box_w / iw, box_h / ih)
        w = int(iw * scale)
        h = int(ih * scale)
        pic.width, pic.height = w, h
        pic.left = int(box_l + (box_w - w) / 2)
        pic.top = int(box_t + (box_h - h) / 2)


def op_set_chart_data(prs, op):
    from pptx.chart.data import CategoryChartData
    _, shape = _shape(prs, op["slide"], op["shape"])
    if not shape.has_chart:
        raise EditError("slide %d shape %d is not a chart" % (op["slide"], op["shape"]))
    chart = shape.chart
    series = op["series"]
    # Default categories to the chart's current ones if not supplied.
    if "categories" in op:
        categories = op["categories"]
    else:
        categories = [str(c) for c in chart.plots[0].categories] if chart.plots else []
    data = CategoryChartData()
    data.categories = categories
    for ser in series:
        data.add_series(ser.get("name", ""), tuple(ser["values"]))
    chart.replace_data(data)
    return "set_chart_data slide %d shape %d (%d series x %d cat)" % (
        op["slide"], op["shape"], len(series), len(categories))


def _picture_px(shape):
    """Pixel (w, h) of a picture shape's embedded image, or None."""
    try:
        return shape.image.size
    except Exception:  # noqa: BLE001
        return None


def _target_box(prs, op):
    """Resolve the box an image op fits into: an existing shape, an inch box,
    or — when neither is given — the full slide (the fixed slide template)."""
    if "into_shape" in op:
        _, target = _shape(prs, op["slide"], op["into_shape"])
        return (target.left, target.top, target.width, target.height)
    if "width_in" in op and "height_in" in op:
        return (Inches(op.get("left_in", 0)), Inches(op.get("top_in", 0)),
                Inches(op["width_in"]), Inches(op["height_in"]))
    return (0, 0, prs.slide_width, prs.slide_height)


def op_fit_image(prs, op):
    """Resize/reposition an existing picture to a fixed box without distortion.

    With no box specified the box is the whole slide, so this snaps an image to
    fill (cover) or fit (contain) the slide template exactly — useful for decks
    whose slides are single full-page images that don't quite line up.
    """
    slide, shape = _shape(prs, op["slide"], op["shape"])
    if shape.shape_type != 13:  # PICTURE
        raise EditError("slide %d shape %d is not a picture" % (op["slide"], op["shape"]))
    fit = op.get("fit", "cover")
    box = _target_box(prs, op)
    _apply_crop(shape, 0, 0, 0, 0)  # clear any prior crop before refitting
    px = _picture_px(shape)
    if px:
        _fit_into_box(shape, px[0], px[1], *box, mode=fit)
    return "fit_image slide %d shape %d (fit=%s)" % (op["slide"], op["shape"], fit)


def op_replace_image(prs, op):
    slide, shape = _shape(prs, op["slide"], op["shape"])
    if shape.shape_type != 13:  # PICTURE
        raise EditError("slide %d shape %d is not a picture" % (op["slide"], op["shape"]))
    path = op["image"]
    if not os.path.exists(path):
        raise EditError("image not found: %s" % path)
    # Remember the existing (template) box so the new image keeps that frame.
    box = (shape.left, shape.top, shape.width, shape.height)
    fit = op.get("fit", "cover")
    # Point the picture's fill at a new image part, keeping its position/size.
    image_part, rId = slide.part.get_or_add_image_part(path)
    shape._element.blipFill.blip.rEmbed = rId
    # Re-fit so a new image with a different aspect ratio isn't distorted.
    _apply_crop(shape, 0, 0, 0, 0)  # clear any prior crop first
    px = _image_px(path)
    if px:
        _fit_into_box(shape, px[0], px[1], *box, mode=fit)
    return "replace_image slide %d shape %d <- %s (fit=%s)" % (
        op["slide"], op["shape"], path, fit)


def op_add_image(prs, op):
    slide = _slide(prs, op["slide"])
    path = op["image"]
    if not os.path.exists(path):
        raise EditError("image not found: %s" % path)

    # Determine the target box. Priority: an existing template shape's geometry,
    # else explicit inch coordinates.
    if "into_shape" in op:
        _, target = _shape(prs, op["slide"], op["into_shape"])
        box = (target.left, target.top, target.width, target.height)
        if op.get("replace_target"):
            target._element.getparent().remove(target._element)
    else:
        box = (Inches(op.get("left_in", 1)), Inches(op.get("top_in", 1)),
               Inches(op["width_in"]) if "width_in" in op else None,
               Inches(op["height_in"]) if "height_in" in op else None)

    fit = op.get("fit", "contain")
    box_l, box_t, box_w, box_h = box
    if box_w is None or box_h is None:
        # No fixed box given: add at native size at the position.
        slide.shapes.add_picture(path, box_l, box_t,
                                 width=box_w, height=box_h)
        return "add_image slide %d <- %s (native size)" % (op["slide"], path)

    pic = slide.shapes.add_picture(path, box_l, box_t, width=box_w, height=box_h)
    px = _image_px(path)
    if px:
        _fit_into_box(pic, px[0], px[1], box_l, box_t, box_w, box_h, mode=fit)
    where = ("into_shape %d" % op["into_shape"]) if "into_shape" in op else "box"
    return "add_image slide %d <- %s (%s, fit=%s)" % (op["slide"], path, where, fit)


def op_delete_shape(prs, op):
    _, shape = _shape(prs, op["slide"], op["shape"])
    shape._element.getparent().remove(shape._element)
    return "delete_shape slide %d shape %d" % (op["slide"], op["shape"])


def op_add_textbox(prs, op):
    slide = _slide(prs, op["slide"])
    tb = slide.shapes.add_textbox(
        Inches(op.get("left_in", 1)), Inches(op.get("top_in", 1)),
        Inches(op.get("width_in", 4)), Inches(op.get("height_in", 1)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if op.get("align") in _ALIGN:
        p.alignment = _ALIGN[op["align"]]
    r = p.add_run()
    r.text = op.get("text", "")
    if "font" in op:
        r.font.name = op["font"]
    if "size_pt" in op:
        r.font.size = Pt(op["size_pt"])
    if "bold" in op:
        r.font.bold = bool(op["bold"])
    if "color" in op:
        r.font.color.rgb = RGBColor.from_string(op["color"].lstrip("#"))
    return "add_textbox slide %d: %r" % (op["slide"], op.get("text", "")[:30])


def _sld_id_list(prs):
    return prs.slides._sldIdLst


def op_delete_slide(prs, op):
    n = op["slide"]
    lst = _sld_id_list(prs)
    ids = list(lst)
    if not (1 <= n <= len(ids)):
        raise EditError("slide %d out of range (deck has %d)" % (n, len(ids)))
    lst.remove(ids[n - 1])
    return "delete_slide %d" % n


def op_move_slide(prs, op):
    frm, to = op["from"], op["to"]
    lst = _sld_id_list(prs)
    ids = list(lst)
    if not (1 <= frm <= len(ids)):
        raise EditError("from %d out of range (deck has %d)" % (frm, len(ids)))
    if not (1 <= to <= len(ids)):
        raise EditError("to %d out of range (deck has %d)" % (to, len(ids)))
    el = ids[frm - 1]
    lst.remove(el)
    lst.insert(to - 1, el)
    return "move_slide %d -> %d" % (frm, to)


OPS = {
    "replace_text": op_replace_text,
    "set_text": op_set_text,
    "set_table_cell": op_set_table_cell,
    "set_chart_data": op_set_chart_data,
    "fit_image": op_fit_image,
    "replace_image": op_replace_image,
    "add_image": op_add_image,
    "add_textbox": op_add_textbox,
    "delete_shape": op_delete_shape,
    "delete_slide": op_delete_slide,
    "move_slide": op_move_slide,
}


def apply_edits(spec, source=None, output=None):
    source = source or spec.get("source")
    output = output or spec.get("output")
    if not source:
        raise EditError("no source .pptx given (spec 'source' or -i)")
    if not output:
        raise EditError("no output path given (spec 'output' or -o)")
    prs = Presentation(source)
    log = []
    for n, op in enumerate(spec.get("operations", []), start=1):
        kind = op.get("op")
        if kind not in OPS:
            raise EditError("operation %d: unknown op %r (valid: %s)"
                            % (n, kind, ", ".join(sorted(OPS))))
        try:
            log.append(OPS[kind](prs, op))
        except EditError:
            raise
        except KeyError as e:
            raise EditError("operation %d (%s): missing field %s" % (n, kind, e))
    prs.save(output)
    return output, log


def main(argv=None):
    ap = argparse.ArgumentParser(description="Apply a JSON edit spec to a .pptx.")
    ap.add_argument("spec", help="path to edit-spec JSON")
    ap.add_argument("-i", "--input", help="source .pptx (overrides spec 'source')")
    ap.add_argument("-o", "--output", help="output .pptx (overrides spec 'output')")
    args = ap.parse_args(argv)

    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)

    try:
        output, log = apply_edits(spec, args.input, args.output)
    except EditError as e:
        print("ERROR: %s" % e, file=sys.stderr)
        return 1
    except Exception as e:  # noqa: BLE001
        print("ERROR: failed to apply edits: %s" % e, file=sys.stderr)
        return 1

    for line in log:
        print("  " + line)
    print("Wrote %s (%d operation(s))" % (output, len(log)))
    return 0


if __name__ == "__main__":
    sys.exit(main())

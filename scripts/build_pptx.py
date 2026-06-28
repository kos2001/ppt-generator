#!/usr/bin/env python3
"""Render a presentation spec (JSON) into an editable .pptx file.

The spec describes WHAT goes on each slide (content + a layout name); the chosen
template describes HOW it looks (palette + fonts). This separation is the whole
point: the same content can be re-rendered in any template, and a deck is built
deterministically from a single JSON file rather than hand-placed shape by shape.

Usage:
    python build_pptx.py spec.json -o out.pptx
    python build_pptx.py spec.json            # writes <title>.pptx next to spec
    echo '{...}' | python build_pptx.py - -o out.pptx   # spec from stdin

The spec schema and layout catalog are documented in references/spec-schema.md
and references/layouts.md. validate_spec.py checks a spec before rendering.
"""
import argparse
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
import math

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from templates import get_theme, DEFAULT_THEME  # noqa: E402
from flowchart_layout import layer_graph, normalize_edges  # noqa: E402

# 16:9 widescreen canvas (13.333" x 7.5"). Most modern decks use this.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)


# --------------------------------------------------------------------------- #
# Low-level drawing helpers
# --------------------------------------------------------------------------- #
def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _blank_slide(prs):
    # slide_layouts[6] is the built-in blank layout: no placeholders, full
    # control. We draw every element ourselves so templates look identical
    # regardless of the underlying PowerPoint theme.
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, theme):
    if theme["background"].upper() == "FFFFFF":
        return  # white is the default; skip the extra shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme["background"])
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Send to back so all later content draws on top.
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_rect(slide, left, top, width, height, fill_hex=None, line_hex=None,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill_hex is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line_hex)
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(slide, text, left, top, width, height, *, font, size,
             color, bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.05, wrap=True,
             shrink=False):
    """Add a text box. `text` may be a string or a list of paragraph dicts.

    A paragraph dict supports: {text, level, bold, size, color, bullet}.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if shrink:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE

    paragraphs = text if isinstance(text, list) else [{"text": text}]
    for i, para in enumerate(paragraphs):
        if isinstance(para, str):
            para = {"text": para}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.level = para.get("level", 0)
        if para.get("space_after") is not None:
            p.space_after = Pt(para["space_after"])
        run = p.add_run()
        run.text = para.get("text", "")
        f = run.font
        f.name = font
        f.size = Pt(para.get("size", size))
        f.bold = para.get("bold", bold)
        f.italic = para.get("italic", italic)
        f.color.rgb = _rgb(para.get("color", color))
    return box


def _slide_title(slide, theme, title, *, eyebrow=None, subtitle=None):
    """Content-slide header. The visual treatment depends on the theme's
    `header_style`:

    - "accent" (default): a small accent block to the left of a colored title
      on the slide background. Light, editorial.
    - "bar": a full-width header band filled with the primary color and a title
      in `on_primary` (e.g. a black header bar with white text). Formal,
      report/document style. No divider line is drawn below it.
    """
    if theme.get("header_style") == "bar":
        _header_bar(slide, theme, title, eyebrow=eyebrow)
        sub_top = Inches(1.12)
    else:
        _header_accent(slide, theme, title, eyebrow=eyebrow)
        sub_top = Inches(1.3)
    # Optional subtitle line directly under the header, above the body. It sits
    # in the gap before any layout's content (which all start at >= 1.7"), so it
    # works uniformly across every content layout without shifting their bodies.
    if subtitle:
        # Intentionally larger than body text so it reads as a lead-in, not a
        # caption. Relative to the theme's body size, overridable via subtitle_pt.
        add_text(slide, subtitle, MARGIN, sub_top, SLIDE_W - 2 * MARGIN,
                 Inches(0.5), font=theme["body_font"],
                 size=theme.get("subtitle_pt", theme["body_pt"] + 4),
                 color=theme["secondary"], italic=True, line_spacing=1.05)


def _header_accent(slide, theme, title, *, eyebrow=None):
    add_rect(slide, MARGIN, Inches(0.62), Inches(0.16), Inches(0.55),
             fill_hex=theme["accent"])
    top = Inches(0.5)
    if eyebrow:
        add_text(slide, eyebrow.upper(), MARGIN + Inches(0.32), Inches(0.42),
                 SLIDE_W - 2 * MARGIN, Inches(0.3),
                 font=theme["body_font"], size=12, color=theme["secondary"],
                 bold=True)
        top = Inches(0.72)
    add_text(slide, title, MARGIN + Inches(0.32), top,
             SLIDE_W - 2 * MARGIN - Inches(0.32), Inches(0.8),
             font=theme["heading_font"], size=theme["heading_pt"],
             color=theme["primary"], bold=True)


def _header_bar(slide, theme, title, *, eyebrow=None):
    bar_h = Inches(1.0)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill_hex=theme["primary"])
    # Optional thin rule directly under the bar (a brand accent line). Themes
    # that don't set `header_rule` (e.g. report) get no divider, as before.
    rule = theme.get("header_rule")
    if rule:
        add_rect(slide, 0, bar_h, SLIDE_W, Inches(0.05), fill_hex=rule)
    # Reserve room at the right for the top-right classification marker (drawn
    # globally per slide by _classification_marker) so a long title never
    # collides with it.
    label_w = Inches(3.0) if theme.get("header_label") else 0
    title_w = SLIDE_W - 2 * MARGIN - label_w
    if eyebrow:
        # Keep the eyebrow + title stacked WITHIN the bar — on a dark bar a title
        # that spills below it would be white text on the white body (invisible).
        add_text(slide, eyebrow.upper(), MARGIN, Inches(0.13),
                 title_w, Inches(0.24),
                 font=theme["body_font"], size=11,
                 color=theme["on_primary"], bold=True)
        add_text(slide, title, MARGIN, Inches(0.38),
                 title_w, Inches(0.54),
                 font=theme["heading_font"], size=theme["heading_pt"],
                 color=theme["on_primary"], bold=True)
    else:
        add_text(slide, title, MARGIN, 0, title_w, bar_h,
                 font=theme["heading_font"], size=theme["heading_pt"],
                 color=theme["on_primary"], bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)


# Layouts whose top-right corner sits on a LIGHT background, so the marker must
# use a dark color to stay legible. Every other layout has a dark area there
# (header bar or full-bleed primary fill) and uses the light on_primary color.
_LIGHT_TOP_LAYOUTS = {"title", "quote", "image"}


def _classification_marker(slide, theme, layout_name, bg_hex=None):
    """Draw the classification label (e.g. "Confidential") in the top-right
    corner of the slide, on every slide, in a color that contrasts with that
    slide's background. Enabled by the theme's `header_label`."""
    label = theme.get("header_label")
    if not label:
        return
    if bg_hex:
        # An explicit per-slide background (e.g. a dark cover) overrides the
        # layout-based guess: pick light/dark text by the background's luminance.
        r, g, b = int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16)
        lum = 0.299 * r + 0.587 * g + 0.114 * b
        color = theme["on_primary"] if lum < 128 else theme["primary"]
    else:
        light_top = layout_name in _LIGHT_TOP_LAYOUTS
        # In bar-style themes the quote slide now carries the dark header bar,
        # so its top is dark — the marker must switch to the light color.
        if light_top and layout_name == "quote" and theme.get("header_style") == "bar":
            light_top = False
        color = theme["primary"] if light_top else theme["on_primary"]
    add_text(slide, label, SLIDE_W - MARGIN - Inches(3.0), Inches(0.2),
             Inches(3.0), Inches(0.6), font=theme["body_font"], size=18,
             color=color, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)


def _footer(slide, theme, slide_no=None, total=None, footer_text=None):
    parts = []
    if footer_text:
        parts.append(footer_text)
    txt = "   ·   ".join(parts) if parts else ""
    if txt:
        add_text(slide, txt, MARGIN, SLIDE_H - Inches(0.5),
                 SLIDE_W - 2 * MARGIN - Inches(1.0), Inches(0.3),
                 font=theme["body_font"], size=10, color=theme["text_muted"])
    if slide_no is not None:
        # Bottom-right page number, forced to "current / total" (e.g. "3 / 15")
        # when the total is known, so the reader always sees their position.
        num = "%d / %d" % (slide_no, total) if total else str(slide_no)
        add_text(slide, num, SLIDE_W - MARGIN - Inches(1.6),
                 SLIDE_H - Inches(0.62), Inches(1.6), Inches(0.42),
                 font=theme["body_font"], size=16, color=theme["text_muted"],
                 align=PP_ALIGN.RIGHT)


def _bullet_paragraphs(theme, items):
    """Turn a list of bullet items into paragraph dicts with bullet glyphs.

    Items may be plain strings or {text, level, bold}. Levels indent and use a
    lighter glyph. We prepend the glyph manually because the blank layout has no
    list formatting of its own.
    """
    glyphs = ["●  ", "–  ", "·  "]
    paras = []
    for item in items:
        if isinstance(item, str):
            item = {"text": item}
        level = min(item.get("level", 0), 2)
        glyph = glyphs[level]
        paras.append({
            "text": glyph + item.get("text", ""),
            "level": level,
            "bold": item.get("bold", False),
            "size": item.get("size", theme["body_pt"] - 2 * level),
            "color": item.get("color", theme["text"] if level == 0 else theme["text_muted"]),
            "space_after": item.get("space_after", 10),
        })
    return paras


# --------------------------------------------------------------------------- #
# Slide layouts — one function per layout name in the spec
# --------------------------------------------------------------------------- #
def layout_title(slide, theme, s):
    # Optional per-slide background override (e.g. a dark cover): fill the whole
    # slide with the given hex and switch the text to light colors so it stays
    # legible. Without it, the title slide renders on the theme background.
    bg = s.get("background")
    if bg:
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_hex=bg)
        title_color = s.get("color", "FFFFFF")
        eyebrow_color = sub_color = meta_color = title_color
    else:
        fill_background(slide, theme)
        title_color = theme["primary"]
        eyebrow_color = theme["secondary"]
        sub_color = meta_color = theme["text_muted"]
    # Big accent block on the left edge for visual weight.
    add_rect(slide, 0, 0, Inches(0.4), SLIDE_H, fill_hex=theme["accent"])
    add_text(slide, s.get("eyebrow", "").upper(), Inches(1.1), Inches(2.1),
             SLIDE_W - Inches(2.0), Inches(0.4),
             font=theme["body_font"], size=14, color=eyebrow_color, bold=True)
    add_text(slide, s.get("title", ""), Inches(1.1), Inches(2.5),
             SLIDE_W - Inches(2.0), Inches(2.0),
             font=theme["heading_font"], size=theme["title_pt"] + 8,
             color=title_color, bold=True, line_spacing=1.0)
    if s.get("subtitle"):
        add_text(slide, s["subtitle"], Inches(1.1), Inches(4.55),
                 SLIDE_W - Inches(2.0), Inches(1.0),
                 font=theme["body_font"], size=theme["body_pt"] + 4,
                 color=sub_color, line_spacing=1.1)
    meta = "   |   ".join(x for x in [s.get("author"), s.get("date")] if x)
    if meta:
        add_text(slide, meta, Inches(1.1), SLIDE_H - Inches(1.0),
                 SLIDE_W - Inches(2.0), Inches(0.4),
                 font=theme["body_font"], size=13, color=meta_color)


def layout_section(slide, theme, s):
    # Full-bleed primary-colored divider slide.
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_hex=theme["primary"])
    if s.get("number"):
        add_text(slide, str(s["number"]), Inches(1.1), Inches(1.8),
                 Inches(4.0), Inches(1.5),
                 font=theme["heading_font"], size=90, color=theme["accent"], bold=True)
    add_text(slide, s.get("title", ""), Inches(1.1), Inches(3.4),
             SLIDE_W - Inches(2.2), Inches(2.0),
             font=theme["heading_font"], size=theme["title_pt"],
             color=theme["on_primary"], bold=True, line_spacing=1.0)
    if s.get("subtitle"):
        add_text(slide, s["subtitle"], Inches(1.1), Inches(5.0),
                 SLIDE_W - Inches(2.2), Inches(1.2),
                 font=theme["body_font"], size=theme["body_pt"] + 2,
                 color=theme["on_primary"], line_spacing=1.15)
    add_rect(slide, Inches(1.1), Inches(3.2), Inches(2.2), Inches(0.08),
             fill_hex=theme["accent"])


def layout_bullets(slide, theme, s):
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    body_top = Inches(1.7)
    if s.get("lead"):
        add_text(slide, s["lead"], MARGIN, body_top,
                 SLIDE_W - 2 * MARGIN, Inches(0.7),
                 font=theme["body_font"], size=theme["body_pt"] + 2,
                 color=theme["secondary"], italic=True, line_spacing=1.1)
        body_top = Inches(2.5)
    bullets = s.get("bullets", [])
    add_text(slide, _bullet_paragraphs(theme, bullets), MARGIN, body_top,
             SLIDE_W - 2 * MARGIN, SLIDE_H - body_top - Inches(0.7),
             font=theme["body_font"], size=theme["body_pt"], color=theme["text"],
             line_spacing=1.1, anchor=MSO_ANCHOR.TOP)


def layout_content(slide, theme, s):
    # Free-form paragraphs (no bullets). `body` is a string or list of strings.
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    body = s.get("body", "")
    if isinstance(body, str):
        body = [body]
    paras = [{"text": p, "size": theme["body_pt"], "color": theme["text"],
              "space_after": 12} for p in body]
    add_text(slide, paras, MARGIN, Inches(1.8),
             SLIDE_W - 2 * MARGIN, SLIDE_H - Inches(2.5),
             font=theme["body_font"], size=theme["body_pt"],
             color=theme["text"], line_spacing=1.25)


def _column_block(slide, theme, col, left, top, width, height):
    """Render one column of a two-up layout. `col` has {heading, bullets|body}."""
    y = top
    if col.get("heading"):
        add_text(slide, col["heading"], left, y, width, Inches(0.5),
                 font=theme["heading_font"], size=theme["body_pt"] + 4,
                 color=theme["secondary"], bold=True)
        y = Emu(y + Inches(0.6))
    if col.get("bullets"):
        add_text(slide, _bullet_paragraphs(theme, col["bullets"]),
                 left, y, width, Emu(top + height - y),
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], line_spacing=1.1)
    elif col.get("body"):
        body = col["body"] if isinstance(col["body"], list) else [col["body"]]
        paras = [{"text": p, "space_after": 10} for p in body]
        add_text(slide, paras, left, y, width, Emu(top + height - y),
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], line_spacing=1.2)


def layout_two_column(slide, theme, s):
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    gap = Inches(0.6)
    col_w = (SLIDE_W - 2 * MARGIN - gap) / 2
    top = Inches(1.9)
    height = SLIDE_H - top - Inches(0.7)
    _column_block(slide, theme, s.get("left", {}), MARGIN, top, col_w, height)
    _column_block(slide, theme, s.get("right", {}),
                  Emu(MARGIN + col_w + gap), top, col_w, height)


def layout_comparison(slide, theme, s):
    # Two contrasting cards side by side (e.g. Before/After, A/B, Pros/Cons).
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    gap = Inches(0.5)
    col_w = (SLIDE_W - 2 * MARGIN - gap) / 2
    top = Inches(1.9)
    height = SLIDE_H - top - Inches(0.7)
    cards = [s.get("left", {}), s.get("right", {})]
    head_fills = [theme["primary"], theme["accent"]]
    for i, card in enumerate(cards):
        left = Emu(MARGIN + i * (col_w + gap))
        add_rect(slide, left, top, col_w, height, fill_hex=theme["surface"])
        add_rect(slide, left, top, col_w, Inches(0.7), fill_hex=head_fills[i])
        add_text(slide, card.get("heading", ""), Emu(left + Inches(0.3)), top,
                 Emu(col_w - Inches(0.6)), Inches(0.7),
                 font=theme["heading_font"], size=theme["body_pt"] + 4,
                 color=theme["on_primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, _bullet_paragraphs(theme, card.get("bullets", [])),
                 Emu(left + Inches(0.3)), Emu(top + Inches(0.95)),
                 Emu(col_w - Inches(0.6)), Emu(height - Inches(1.2)),
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], line_spacing=1.1)


def layout_metrics(slide, theme, s):
    # Up to 4 KPI cards in a row.
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    metrics = s.get("metrics", [])[:4]
    n = max(len(metrics), 1)
    gap = Inches(0.4)
    total_w = SLIDE_W - 2 * MARGIN
    card_w = (total_w - gap * (n - 1)) / n
    top = Inches(2.3)
    card_h = Inches(2.8)
    for i, m in enumerate(metrics):
        left = Emu(MARGIN + i * (card_w + gap))
        add_rect(slide, left, top, card_w, card_h, fill_hex=theme["surface"])
        add_rect(slide, left, top, card_w, Inches(0.12), fill_hex=theme["accent"])
        add_text(slide, str(m.get("value", "")), left, Emu(top + Inches(0.5)),
                 card_w, Inches(1.1),
                 font=theme["heading_font"], size=48, color=theme["primary"],
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, m.get("label", ""), left, Emu(top + Inches(1.65)),
                 card_w, Inches(0.5),
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], bold=True, align=PP_ALIGN.CENTER)
        if m.get("sublabel"):
            add_text(slide, m["sublabel"], Emu(left + Inches(0.2)),
                     Emu(top + Inches(2.1)), Emu(card_w - Inches(0.4)),
                     Inches(0.6), font=theme["body_font"], size=12,
                     color=theme["text_muted"], align=PP_ALIGN.CENTER,
                     line_spacing=1.0)


def layout_quote(slide, theme, s):
    fill_background(slide, theme)
    # Bar-style themes (e.g. samsung/report) carry a header bar on every content
    # slide; draw it here too so the quote slide matches the rest of the deck
    # instead of floating on a bare background. The quote then sits below the bar.
    if theme.get("header_style") == "bar":
        _header_bar(slide, theme, s.get("title", ""))
        qm_top, q_top, attr_top = Inches(1.6), Inches(2.8), Inches(5.9)
    else:
        qm_top, q_top, attr_top = Inches(0.9), Inches(2.3), Inches(5.6)
    add_text(slide, "“", Inches(0.8), qm_top, Inches(3.0), Inches(2.0),
             font=theme["heading_font"], size=160, color=theme["accent"], bold=True)
    add_text(slide, s.get("quote", ""), Inches(1.6), q_top,
             SLIDE_W - Inches(3.2), Inches(2.8),
             font=theme["heading_font"], size=theme["heading_pt"] + 2,
             color=theme["primary"], italic=True, line_spacing=1.2,
             anchor=MSO_ANCHOR.MIDDLE)
    if s.get("attribution"):
        add_text(slide, "— " + s["attribution"], Inches(1.6), attr_top,
                 SLIDE_W - Inches(3.2), Inches(0.6),
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text_muted"], bold=True)


def _image_dims(path):
    try:
        from PIL import Image
        with Image.open(path) as im:
            return im.size  # (w, h) in px
    except Exception:
        return None


def _place_image(slide, path, box_left, box_top, box_w, box_h):
    """Place an image fitted (contain) inside the given box, centered."""
    dims = _image_dims(path)
    if not dims:
        # Missing/unreadable image: draw a labeled placeholder instead of crashing.
        add_rect(slide, box_left, box_top, box_w, box_h, fill_hex="DDDDDD")
        add_text(slide, "[image not found:\n%s]" % os.path.basename(path),
                 box_left, box_top, box_w, box_h, font="Arial", size=12,
                 color="888888", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    iw, ih = dims
    scale = min(box_w / iw, box_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    left = int(box_left + (box_w - w) / 2)
    top = int(box_top + (box_h - h) / 2)
    slide.shapes.add_picture(path, left, top, width=w, height=h)


def layout_image(slide, theme, s):
    """Image with optional text. `position`: right (default), left, or full."""
    fill_background(slide, theme)
    pos = s.get("position", "right")
    img = s.get("image")
    if pos == "full":
        if img:
            _place_image(slide, img, 0, 0, SLIDE_W, SLIDE_H)
        if s.get("title"):
            add_rect(slide, 0, SLIDE_H - Inches(1.3), SLIDE_W, Inches(1.3),
                     fill_hex=theme["primary"])
            add_text(slide, s["title"], MARGIN, SLIDE_H - Inches(1.15),
                     SLIDE_W - 2 * MARGIN, Inches(1.0),
                     font=theme["heading_font"], size=theme["heading_pt"],
                     color=theme["on_primary"], bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        return
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    half = (SLIDE_W - 2 * MARGIN - Inches(0.5)) / 2
    top = Inches(1.9)
    height = SLIDE_H - top - Inches(0.7)
    if pos == "left":
        img_left, txt_left = MARGIN, Emu(MARGIN + half + Inches(0.5))
    else:
        txt_left, img_left = MARGIN, Emu(MARGIN + half + Inches(0.5))
    if img:
        _place_image(slide, img, img_left, top, half, height)
    bullets = s.get("bullets")
    if bullets:
        add_text(slide, _bullet_paragraphs(theme, bullets), txt_left, top,
                 half, height, font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], line_spacing=1.15)
    elif s.get("body"):
        body = s["body"] if isinstance(s["body"], list) else [s["body"]]
        paras = [{"text": p, "space_after": 10} for p in body]
        add_text(slide, paras, txt_left, top, half, height,
                 font=theme["body_font"], size=theme["body_pt"],
                 color=theme["text"], line_spacing=1.25)


def layout_table(slide, theme, s):
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    tbl_spec = s.get("table", {})
    header = tbl_spec.get("header", [])
    rows = tbl_spec.get("rows", [])
    ncols = max(len(header), max((len(r) for r in rows), default=1))
    nrows = len(rows) + (1 if header else 0)
    if nrows == 0 or ncols == 0:
        return
    top = Inches(2.0)
    width = SLIDE_W - 2 * MARGIN
    height = min(SLIDE_H - top - Inches(0.7), Inches(0.5) * nrows + Inches(0.3))
    gframe = slide.shapes.add_table(nrows, ncols, MARGIN, top, width, height)
    table = gframe.table
    r0 = 0
    if header:
        for c in range(ncols):
            cell = table.cell(0, c)
            cell.text = str(header[c]) if c < len(header) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(theme["primary"])
            _style_cell(cell, theme, theme["on_primary"], bold=True)
        r0 = 1
    for ri, row in enumerate(rows):
        for c in range(ncols):
            cell = table.cell(ri + r0, c)
            cell.text = str(row[c]) if c < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                theme["background"] if ri % 2 == 0 else theme["surface"])
            _style_cell(cell, theme, theme["text"])


def _style_cell(cell, theme, color_hex, bold=False):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(theme["body_pt"] - 3)
            run.font.name = theme["body_font"]
            run.font.bold = bold
            run.font.color.rgb = _rgb(color_hex)


_CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
}


def layout_chart(slide, theme, s):
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    spec = s.get("chart", {})
    ctype = _CHART_TYPES.get(spec.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = spec.get("categories", [])
    series = spec.get("series", [])
    if series and isinstance(series[0], dict):
        for ser in series:
            data.add_series(ser.get("name", ""), tuple(ser.get("values", [])))
    else:
        # Bare list of numbers => single unnamed series.
        data.add_series(spec.get("series_name", "Series 1"), tuple(series))
    top = Inches(1.9)
    gframe = slide.shapes.add_chart(
        ctype, MARGIN, top, SLIDE_W - 2 * MARGIN, SLIDE_H - top - Inches(0.7),
        data)
    chart = gframe.chart
    chart.has_title = False
    multi = len(series) > 1 and isinstance(series[0], dict)
    if multi or spec.get("type") == "pie":
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    try:
        for plot in chart.plots:
            plot.has_data_labels = bool(spec.get("data_labels", False))
    except Exception:
        pass


def layout_closing(slide, theme, s):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_hex=theme["primary"])
    add_rect(slide, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4),
             fill_hex=theme["accent"])
    add_text(slide, s.get("title", "Thank You"), MARGIN, Inches(2.6),
             SLIDE_W - 2 * MARGIN, Inches(1.6),
             font=theme["heading_font"], size=theme["title_pt"] + 6,
             color=theme["on_primary"], bold=True, align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        add_text(slide, s["subtitle"], MARGIN, Inches(4.3),
                 SLIDE_W - 2 * MARGIN, Inches(1.2),
                 font=theme["body_font"], size=theme["body_pt"] + 2,
                 color=theme["on_primary"], align=PP_ALIGN.CENTER,
                 line_spacing=1.2)


# --------------------------------------------------------------------------- #
# Diagrams — drawn as NATIVE PowerPoint shapes (autoshapes, connectors,
# freeforms), never flattened images, so every box, arrow, and label stays
# fully editable in PowerPoint/Keynote after the deck is generated.
# --------------------------------------------------------------------------- #
def _autoshape(slide, shape_type, left, top, w, h, fill, line=None, line_w=1.25):
    sp = slide.shapes.add_shape(shape_type, int(left), int(top), int(w), int(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
    if line:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _shape_label(shape, theme, title, desc=None, color=None,
                 title_size=14, desc_size=11, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for attr in ("margin_left", "margin_right"):
        setattr(tf, attr, Inches(0.08))
    for attr in ("margin_top", "margin_bottom"):
        setattr(tf, attr, Inches(0.03))
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(title_size)
    r.font.name = theme["heading_font"]
    r.font.color.rgb = _rgb(color or theme["on_primary"])
    if desc:
        p2 = tf.add_paragraph()
        p2.alignment = align
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(desc_size)
        r2.font.name = theme["body_font"]
        r2.font.color.rgb = _rgb(color or theme["on_primary"])
    return shape


def _connect(slide, x1, y1, x2, y2, color, width=1.75, arrow=True, dash=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   int(x1), int(y1), int(x2), int(y2))
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dash:
        # OOXML: a:prstDash must precede a:tailEnd in the line element, and the
        # color set above already added a:solidFill before this — so appending
        # here keeps the schema's child order valid.
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if arrow:
        ln.append(ln.makeelement(
            qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    c.shadow.inherit = False
    return c


def _freeform_poly(slide, pts, fill, line=None):
    ff = slide.shapes.build_freeform(float(pts[0][0]), float(pts[0][1]), scale=1.0)
    ff.add_line_segments([(float(x), float(y)) for x, y in pts[1:]], close=True)
    shp = ff.convert_to_shape()
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    if line:
        shp.line.color.rgb = _rgb(line)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _diag_area():
    """The drawing rectangle below the slide header: (left, top, width, height)."""
    top = Inches(2.0)
    return MARGIN, top, SLIDE_W - 2 * MARGIN, SLIDE_H - top - Inches(0.7)


def _band_colors(theme, n):
    # Alternate primary/secondary so adjacent shapes read as distinct.
    pair = [theme["primary"], theme["secondary"]]
    return [pair[i % 2] for i in range(n)]


def _diag_process(slide, theme, nodes):
    left, top, w, h = _diag_area()
    n = max(len(nodes), 1)
    arrow_w = Inches(0.5)
    box_w = (w - arrow_w * (n - 1)) / n
    box_h = Inches(1.7)
    box_top = top + (h - box_h) / 2
    colors = _band_colors(theme, n)
    for i, node in enumerate(nodes):
        x = left + i * (box_w + arrow_w)
        sp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h,
                        colors[i])
        _shape_label(sp, theme, node.get("title", ""), node.get("desc"),
                     title_size=15, desc_size=11)
        if i < n - 1:
            ax = x + box_w
            ay = box_top + box_h / 2
            _autoshape(slide, MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.05),
                       ay - Inches(0.18), arrow_w - Inches(0.1), Inches(0.36),
                       theme["accent"])


def _diag_cycle(slide, theme, nodes):
    left, top, w, h = _diag_area()
    n = max(len(nodes), 1)
    cx, cy = left + w / 2, top + h / 2
    box_w, box_h = Inches(2.3), Inches(1.05)
    radius = min(w, h) / 2 - box_h
    colors = _band_colors(theme, n)
    centers = []
    for i in range(n):
        ang = -math.pi / 2 + i * 2 * math.pi / n  # start at top, clockwise
        px = cx + radius * math.cos(ang)
        py = cy + radius * math.sin(ang)
        centers.append((px, py))
    # Arrows first (so boxes sit on top), connecting node i -> i+1 around the
    # ring. Draw each from the edge of one box to the edge of the next (where the
    # center-to-center line crosses each box's rectangle) so the arrow spans the
    # gap and touches both boxes, instead of floating as a stub in the middle.
    hw, hh = box_w / 2, box_h / 2

    def _edge(px, py, ux, uy):
        tx = hw / abs(ux) if ux else float("inf")
        ty = hh / abs(uy) if uy else float("inf")
        t = min(tx, ty)
        return px + ux * t, py + uy * t

    for i in range(n):
        x1, y1 = centers[i]
        x2, y2 = centers[(i + 1) % n]
        dx, dy = x2 - x1, y2 - y1
        d = math.hypot(dx, dy) or 1
        ux, uy = dx / d, dy / d
        sx, sy = _edge(x1, y1, ux, uy)
        ex, ey = _edge(x2, y2, -ux, -uy)
        _connect(slide, sx, sy, ex, ey, theme["accent"], width=2.0)
    for i, node in enumerate(nodes):
        px, py = centers[i]
        sp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        px - box_w / 2, py - box_h / 2, box_w, box_h, colors[i])
        _shape_label(sp, theme, node.get("title", ""), node.get("desc"),
                     title_size=13, desc_size=10)


def _diag_hierarchy(slide, theme, nodes):
    # Two levels: nodes[0] is the root, the rest are its children.
    left, top, w, h = _diag_area()
    if not nodes:
        return
    root, children = nodes[0], nodes[1:]
    root_w, root_h = Inches(3.0), Inches(1.0)
    root_x = left + (w - root_w) / 2
    rsp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, root_x, top, root_w, root_h,
                     theme["primary"])
    _shape_label(rsp, theme, root.get("title", ""), root.get("desc"),
                 title_size=16, desc_size=11)
    if not children:
        return
    m = len(children)
    gap = Inches(0.4)
    child_w = min((w - gap * (m - 1)) / m, Inches(3.2))
    total = child_w * m + gap * (m - 1)
    cstart = left + (w - total) / 2
    child_h = Inches(1.1)
    child_top = top + Inches(2.4)
    rcx, rcy = root_x + root_w / 2, top + root_h
    for i, child in enumerate(children):
        cx = cstart + i * (child_w + gap)
        ccx = cx + child_w / 2
        _connect(slide, rcx, rcy + Inches(0.1), ccx, child_top, theme["secondary"],
                 width=1.5, arrow=False)
        csp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, child_top,
                         child_w, child_h, theme["surface"], line=theme["secondary"])
        _shape_label(csp, theme, child.get("title", ""), child.get("desc"),
                     color=theme["text"], title_size=13, desc_size=10)


def _diag_pyramid(slide, theme, nodes, funnel=False):
    left, top, w, h = _diag_area()
    m = max(len(nodes), 1)
    base_hw = (w * 0.72) / 2
    cx = left + w / 2
    level_h = h / m
    colors = _band_colors(theme, m)

    def hw(depth):  # half-width at fractional depth 0..m (0 = apex)
        frac = depth / m
        return base_hw * (frac if not funnel else (1 - frac))

    for i, node in enumerate(nodes):
        y_top = top + i * level_h
        y_bot = top + (i + 1) * level_h
        hwt, hwb = hw(i), hw(i + 1)
        pts = [(cx - hwt, y_top), (cx + hwt, y_top),
               (cx + hwb, y_bot), (cx - hwb, y_bot)]
        _freeform_poly(slide, pts, colors[i])
        # Label centered on the band (its own textbox — band polygon stays clean).
        lbl = node.get("title", "")
        if node.get("desc"):
            lbl += "  —  " + node["desc"]
        add_text(slide, lbl, left, y_top, w, level_h,
                 font=theme["heading_font"], size=14, color=theme["on_primary"],
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 wrap=True)


def _diag_funnel(slide, theme, nodes):
    _diag_pyramid(slide, theme, nodes, funnel=True)


def _diag_timeline(slide, theme, nodes):
    left, top, w, h = _diag_area()
    n = max(len(nodes), 1)
    line_y = top + h / 2
    pad = Inches(0.4)
    x0, x1 = left + pad, left + w - pad
    # The spine.
    _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x0, line_y - Inches(0.03),
               x1 - x0, Inches(0.06), theme["secondary"])
    step = (x1 - x0) / (n - 1) if n > 1 else 0
    colors = _band_colors(theme, n)
    box_w = Inches(2.4)
    for i, node in enumerate(nodes):
        mx = x0 if n == 1 else x0 + i * step
        # Marker dot on the spine.
        r = Inches(0.13)
        _autoshape(slide, MSO_SHAPE.OVAL, mx - r, line_y - r, 2 * r, 2 * r,
                   colors[i])
        above = (i % 2 == 0)
        by = (line_y - Inches(1.55)) if above else (line_y + Inches(0.35))
        bx = max(left, min(mx - box_w / 2, left + w - box_w))
        sp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, box_w,
                        Inches(1.2), theme["surface"], line=theme["secondary"])
        _shape_label(sp, theme, node.get("title", ""), node.get("desc"),
                     color=theme["text"], title_size=13, desc_size=10)
        # Stem from marker to the label box.
        sy = (by + Inches(1.2)) if above else by
        _connect(slide, mx, line_y, mx, sy, theme["secondary"], width=1.25,
                 arrow=False)


def _diag_flowchart(slide, theme, spec):
    """Auto-laid-out flowchart from nodes + edges, drawn as native shapes.

    The "Mermaid-like" capability that stays editable: a layered (Sugiyama-style)
    layout positions an arbitrary directed graph — branches, merges, and
    loop-backs — then draws real rounded-rectangle boxes and arrow connectors.
    Unlike a Mermaid image, every box stays movable/re-typable.

    spec = {nodes: [{id, title, desc}],
            edges: [[from_id, to_id] | {from, to, label}],
            direction: "LR" (default) | "TD"}
    Layering and back-edge detection are shared (flowchart_layout.layer_graph);
    feedback edges are drawn as muted return paths so cycles stay compact.
    """
    nodes = spec.get("nodes", [])
    if not nodes:
        return
    td = str(spec.get("direction", "LR")).upper() == "TD"
    ids, nodemap = [], {}
    for i, nd in enumerate(nodes):
        nid = str(nd.get("id", i))
        ids.append(nid)
        nodemap[nid] = nd
    edges = [(a, b, l) for (a, b, l) in normalize_edges(spec.get("edges", []))
             if a in nodemap and b in nodemap]
    rank, layers, feedback = layer_graph(ids, edges)
    n_layers = max(layers) + 1
    max_in_layer = max(len(v) for v in layers.values())

    left, top, w, h = _diag_area()
    gap_x, gap_y = Inches(0.55), Inches(0.4)
    if td:  # layers stack top->down; siblings spread across each row
        box_w = min(Inches(2.6), (w - gap_x * (max_in_layer - 1)) / max_in_layer)
        box_h = min(Inches(1.05), (h - gap_y * (n_layers - 1)) / n_layers)
    else:   # layers march left->right; siblings stack within each column
        box_w = min(Inches(2.3), (w - gap_x * (n_layers - 1)) / n_layers)
        box_h = min(Inches(1.2), (h - gap_y * (max_in_layer - 1)) / max_in_layer)

    rects = {}
    for lvl in range(n_layers):
        group = layers.get(lvl, [])
        m = len(group)
        if td:
            row_w = m * box_w + (m - 1) * gap_x
            x0 = left + (w - row_w) / 2
            y = top + lvl * (box_h + gap_y)
            for j, n in enumerate(group):
                rects[n] = (x0 + j * (box_w + gap_x), y)
        else:
            col_h = m * box_h + (m - 1) * gap_y
            y0 = top + (h - col_h) / 2
            x = left + lvl * (box_w + gap_x)
            for j, n in enumerate(group):
                rects[n] = (x, y0 + j * (box_h + gap_y))

    def _mid(n):
        x, y = rects[n]
        return x + box_w / 2, y + box_h / 2

    # Connectors first so boxes sit on top.
    clearance = Inches(0.22)
    for a, b, label in edges:
        if a == b:
            continue  # self-loop: skip (rare); the box itself implies it
        ax, ay = rects[a]
        bx, by = rects[b]
        if (a, b) not in feedback:
            # Forward step: a straight solid arrow between adjacent box edges.
            color = theme["accent"]
            if td:
                p1 = (ax + box_w / 2, ay + box_h)   # bottom-mid
                p2 = (bx + box_w / 2, by)           # top-mid
            else:
                p1 = (ax + box_w, ay + box_h / 2)   # right-mid
                p2 = (bx, by + box_h / 2)           # left-mid
            _connect(slide, p1[0], p1[1], p2[0], p2[1], color, width=1.75)
            lx, ly = (p1[0] + p2[0]) / 2, (p1[1] + p2[1]) / 2
        else:
            # Feedback / loop-back: route a DASHED orthogonal path AROUND the
            # boxes (over the top for LR, out the left for TD), with the
            # arrowhead only on the final leg landing on the target box. A
            # straight diagonal here would overlap the forward arrows and read
            # as pointing at the wrong box; routing clear of them and dashing it
            # makes the return direction unmistakable.
            col = theme["secondary"]
            if not td:  # LR flow -> out the right of the source, up, over, down
                yline = top - clearance        # above every box, below the header
                x_out = min(ax + box_w + clearance, left + w)  # right of the column
                xb, ya = bx + box_w / 2, ay + box_h / 2
                _connect(slide, ax + box_w, ya, x_out, ya, col, 1.5, arrow=False, dash=True)
                _connect(slide, x_out, ya, x_out, yline, col, 1.5, arrow=False, dash=True)
                _connect(slide, x_out, yline, xb, yline, col, 1.5, arrow=False, dash=True)
                _connect(slide, xb, yline, xb, by, col, 1.5, arrow=True, dash=True)
                lx, ly = (x_out + xb) / 2, yline - Inches(0.02)
            else:       # TD flow -> route out the left side
                xline = left - clearance if left - clearance > Inches(0.1) else Inches(0.1)
                ya, yb = ay + box_h / 2, by + box_h / 2
                _connect(slide, ax, ya, xline, ya, col, 1.5, arrow=False, dash=True)
                _connect(slide, xline, ya, xline, yb, col, 1.5, arrow=False, dash=True)
                _connect(slide, xline, yb, bx, yb, col, 1.5, arrow=True, dash=True)
                lx, ly = xline, (ya + yb) / 2
        if label:
            add_text(slide, str(label), lx - Inches(0.7), ly - Inches(0.22),
                     Inches(1.4), Inches(0.3), font=theme["body_font"], size=10,
                     color=theme["text_muted"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    colors = _band_colors(theme, n_layers)
    for n in ids:
        x, y = rects[n]
        sp = _autoshape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h,
                        colors[rank[n]])
        nd = nodemap[n]
        _shape_label(sp, theme, nd.get("title", n), nd.get("desc"),
                     title_size=13, desc_size=10)


_DIAGRAMS = {
    "process": _diag_process,
    "cycle": _diag_cycle,
    "hierarchy": _diag_hierarchy,
    "pyramid": _diag_pyramid,
    "funnel": _diag_funnel,
    "timeline": _diag_timeline,
    "flowchart": _diag_flowchart,
}


def layout_diagram(slide, theme, s):
    fill_background(slide, theme)
    _slide_title(slide, theme, s.get("title", ""), eyebrow=s.get("eyebrow"),
                 subtitle=s.get("subtitle"))
    spec = s.get("diagram", {})
    dtype = spec.get("type", "process")
    nodes = spec.get("nodes", [])
    fn = _DIAGRAMS.get(dtype, _diag_process)
    if dtype not in _DIAGRAMS:
        print("WARN: unknown diagram type '%s' -> using 'process'" % dtype,
              file=sys.stderr)
    # flowchart needs the edges too, so it takes the whole diagram spec; the
    # other types are positional patterns that only need the node list.
    if dtype == "flowchart":
        _diag_flowchart(slide, theme, spec)
    else:
        fn(slide, theme, nodes)


LAYOUTS = {
    "title": layout_title,
    "diagram": layout_diagram,
    "section": layout_section,
    "bullets": layout_bullets,
    "content": layout_content,
    "two_column": layout_two_column,
    "comparison": layout_comparison,
    "metrics": layout_metrics,
    "quote": layout_quote,
    "image": layout_image,
    "table": layout_table,
    "chart": layout_chart,
    "closing": layout_closing,
}


# --------------------------------------------------------------------------- #
# Driver
# --------------------------------------------------------------------------- #
def build(spec, out_path):
    theme_name = spec.get("template", DEFAULT_THEME)
    theme = dict(get_theme(theme_name))  # copy: we may inject per-deck overrides
    if theme_name.strip().lower() not in __import__("templates").THEMES:
        print("WARN: unknown template '%s' -> using '%s'"
              % (theme_name, DEFAULT_THEME), file=sys.stderr)
    # Per-deck override of the header classification label ("" disables it).
    if "header_label" in spec:
        theme["header_label"] = spec["header_label"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = spec.get("slides", [])
    footer_default = spec.get("footer")
    show_numbers = spec.get("slide_numbers", True)

    for i, s in enumerate(slides):
        name = s.get("layout", "bullets")
        fn = LAYOUTS.get(name)
        slide = _blank_slide(prs)
        if fn is None:
            print("WARN: unknown layout '%s' on slide %d -> using 'bullets'"
                  % (name, i + 1), file=sys.stderr)
            fn = layout_bullets
        fn(slide, theme, s)
        # Footers/numbers only on interior content slides, not full-bleed ones.
        if name not in ("title", "section", "closing", "image"):
            # Themes can opt out of footer text (e.g. samsung) so only the
            # page number remains on the right.
            ftext = s.get("footer", footer_default)
            if not theme.get("show_footer_text", True):
                ftext = None
            _footer(slide, theme,
                    slide_no=(i + 1) if show_numbers else None,
                    total=len(slides),
                    footer_text=ftext)
        # Classification marker (e.g. Confidential) on every slide's top-right.
        _classification_marker(slide, theme, name, bg_hex=s.get("background"))

    prs.save(out_path)
    return len(slides)


def _default_out(spec, spec_path):
    title = spec.get("title") or "presentation"
    slug = re.sub(r"[^A-Za-z0-9._-]+", "_", title).strip("_") or "presentation"
    base = os.path.dirname(os.path.abspath(spec_path)) if spec_path and spec_path != "-" else os.getcwd()
    return os.path.join(base, slug + ".pptx")


def main(argv=None):
    ap = argparse.ArgumentParser(description="Render a presentation spec JSON into a .pptx")
    ap.add_argument("spec", help="Path to spec JSON, or '-' for stdin")
    ap.add_argument("-o", "--output", help="Output .pptx path")
    args = ap.parse_args(argv)

    if args.spec == "-":
        spec = json.load(sys.stdin)
    else:
        with open(args.spec, encoding="utf-8") as f:
            spec = json.load(f)

    out = args.output or _default_out(spec, args.spec)
    n = build(spec, out)
    print("Wrote %d slides -> %s" % (n, out))
    return 0


if __name__ == "__main__":
    sys.exit(main())

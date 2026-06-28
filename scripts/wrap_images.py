#!/usr/bin/env python3
"""Wrap an image-only deck (or a folder of slide images) in a template frame.

Image-only decks — slides exported as one full-page picture each (NotebookLM,
translated, or design-tool exports) — have no editable text. A common need is
to keep the images but put them inside a house template's chrome: a branded
header bar, logo/classification marker, and page numbers. Doing that by hand
means rendering a chrome-only frame deck and fitting each image into its body,
once per image — this script does exactly that in one step.

It does NOT make the text editable (the content stays baked into the pixels).
To get editable text instead, transcribe the images into a presentation spec
and render with build_pptx.py (the generation path). See SKILL.md.

Usage:
    # from an existing image-only .pptx (one picture per slide):
    python wrap_images.py --from-pptx deck.pptx -o out.pptx
    # from a folder of images (natural-sorted: slide1, slide2, … slide10):
    python wrap_images.py --images ./imgs --template samsung -o out.pptx
    # options:
    #   --template NAME   template/theme (default: renderer default)
    #   --eyebrow TEXT    running label shown in the header bar
    #   --footer TEXT     footer text on each slide
    #   --fit MODE        contain (default) | cover | stretch
    #   --no-numbers      hide page numbers

The output slide count equals the number of images, so the original deck's
length is preserved.
"""
import argparse
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import build_pptx
import edit_pptx
from templates import get_theme, body_box, DEFAULT_THEME

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def _natural_key(name):
    """Sort 'slide2' before 'slide10' by splitting digit runs into ints."""
    return [int(t) if t.isdigit() else t.lower()
            for t in re.split(r"(\d+)", name)]


def images_from_folder(folder):
    exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
    files = [f for f in os.listdir(folder)
             if os.path.splitext(f)[1].lower() in exts]
    files.sort(key=_natural_key)
    return [os.path.join(folder, f) for f in files]


def images_from_pptx(pptx_path, out_dir):
    """Extract the largest picture from each slide, in slide order, as the
    slide's image. Returns the list of written file paths (one per slide that
    has a picture)."""
    from pptx import Presentation
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    paths = []
    for i, slide in enumerate(prs.slides, start=1):
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # PICTURE
        if not pics:
            continue
        # the slide's main image = the largest picture by area
        pic = max(pics, key=lambda s: (s.width or 0) * (s.height or 0))
        img = pic.image
        ext = img.ext or "png"
        path = os.path.join(out_dir, "slide%d.%s" % (i, ext))
        with open(path, "wb") as f:
            f.write(img.blob)
        paths.append(path)
    return paths


def wrap(images, template, output, *, eyebrow=None, footer=None,
         fit="contain", numbers=True):
    if not images:
        raise SystemExit("no images to wrap")
    theme = get_theme(template)

    # 1) Render a chrome-only frame deck: one empty-body content slide per image.
    frame_spec = {
        "template": template or DEFAULT_THEME,
        "title": os.path.splitext(os.path.basename(output))[0],
        "slide_numbers": numbers,
        "slides": [{"layout": "bullets",
                    "eyebrow": eyebrow or "",
                    "title": "",
                    "bullets": [],
                    **({"footer": footer} if footer else {})}
                   for _ in images],
    }
    frame_path = output + ".frame.tmp.pptx"
    build_pptx.build(frame_spec, frame_path)

    # 2) Fit each image into the body box of its slide.
    left, top, width, height = body_box(theme)
    ops = [{"op": "add_image", "slide": i, "image": img,
            "left_in": left, "top_in": top, "width_in": width,
            "height_in": height, "fit": fit}
           for i, img in enumerate(images, start=1)]
    edit_pptx.apply_edits({"source": frame_path, "output": output,
                           "operations": ops})
    os.remove(frame_path)
    return len(images)


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    src = ap.add_mutually_exclusive_group(required=True)
    src.add_argument("--images", help="folder of slide images")
    src.add_argument("--from-pptx", help="image-only .pptx to pull images from")
    ap.add_argument("-o", "--output", required=True, help="output .pptx")
    ap.add_argument("--template", default=DEFAULT_THEME, help="template/theme name")
    ap.add_argument("--eyebrow", help="running label shown in the header bar")
    ap.add_argument("--footer", help="footer text on each slide (convention: "
                                      "leave off; footer shows only the page number)")
    ap.add_argument("--fit", default="contain", choices=("contain", "cover", "stretch"))
    ap.add_argument("--no-numbers", action="store_true", help="hide page numbers")
    args = ap.parse_args(argv)

    if args.images:
        images = images_from_folder(args.images)
    else:
        images = images_from_pptx(args.from_pptx, args.output + ".imgs")

    n = wrap(images, args.template, args.output,
             eyebrow=args.eyebrow, footer=args.footer,
             fit=args.fit, numbers=not args.no_numbers)
    print("Wrote %s (%d image slide(s), template=%s)" % (args.output, n, args.template))
    return 0


if __name__ == "__main__":
    sys.exit(main())

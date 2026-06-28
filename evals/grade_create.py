#!/usr/bin/env python3
"""Feature extractor / grader for Mode A (create) eval outputs.

Opens a generated .pptx and emits objective features the create-mode assertions
lean on: validity, slide count, which template was used (inferred from heading
font + accent color), native-chart vs picture counts, table count, and per-slide
titles. Pure python-pptx — runs anywhere, no COM.

Template inference matters because the skill's standing default is samsung; an
eval that explicitly asks for, say, the dark template must actually produce dark,
not silently fall back to the default. We match the deck's dominant heading font
and accent fill against scripts/templates.py THEMES.

Usage:
    python evals/grade_create.py <deck.pptx>
    python evals/grade_create.py <deck.pptx> --expect-template dark --min 10 --max 14
"""
import argparse
import json
import sys
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, "scripts")
try:
    from templates import THEMES
except Exception:  # pragma: no cover - templates optional
    THEMES = {}


def _hexrgb(color):
    try:
        return str(color.rgb)
    except Exception:
        return None


def extract(path):
    prs = Presentation(path)
    fonts = Counter()
    fills = Counter()
    n_chart = n_table = n_pic = 0
    titles = []
    for slide in prs.slides:
        slide_title = None
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_pic += 1
            if getattr(sh, "has_chart", False) and sh.has_chart:
                n_chart += 1
            if getattr(sh, "has_table", False) and sh.has_table:
                n_table += 1
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if txt and slide_title is None:
                    slide_title = txt.splitlines()[0][:60]
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
            try:
                if sh.fill.type is not None:
                    c = _hexrgb(sh.fill.fore_color)
                    if c:
                        fills[c] += 1
            except Exception:
                pass
        titles.append(slide_title or "")

    # infer template: score each theme by heading-font match + accent presence
    top_font = fonts.most_common(1)[0][0] if fonts else None
    guess, best = None, -1
    for name, t in THEMES.items():
        score = 0
        if t.get("heading_font") and any(t["heading_font"] == f for f in fonts):
            score += 2
        if t.get("body_font") and any(t["body_font"] == f for f in fonts):
            score += 1
        if t.get("accent") and t["accent"] in fills:
            score += 2
        if t.get("background") and t["background"] in fills:
            score += 1
        if score > best:
            guess, best = name, score

    return {
        "valid": True,
        "slides": len(prs.slides),
        "charts": n_chart,
        "tables": n_table,
        "pictures": n_pic,
        "top_font": top_font,
        "fonts": dict(fonts.most_common(5)),
        "guessed_template": guess,
        "titles": titles,
    }


def main(argv=None):
    ap = argparse.ArgumentParser(description="Extract features from a create-mode deck.")
    ap.add_argument("pptx")
    ap.add_argument("--expect-template")
    ap.add_argument("--min", type=int)
    ap.add_argument("--max", type=int)
    args = ap.parse_args(argv)

    try:
        feats = extract(args.pptx)
    except Exception as e:  # noqa: BLE001
        print(json.dumps({"valid": False, "error": str(e)}, ensure_ascii=False))
        return 1

    checks = []
    if args.expect_template:
        ok = feats["guessed_template"] == args.expect_template
        checks.append(("template == %s" % args.expect_template, ok,
                       "guessed %s (fonts %s)" % (feats["guessed_template"], feats["fonts"])))
    if args.min is not None and args.max is not None:
        ok = args.min <= feats["slides"] <= args.max
        checks.append(("slide count in [%d,%d]" % (args.min, args.max), ok,
                       "%d slides" % feats["slides"]))

    for text, ok, ev in checks:
        print("[%s] %s — %s" % ("PASS" if ok else "FAIL", text, ev))
    print(json.dumps(feats, ensure_ascii=False))
    return 0 if all(ok for _, ok, _ in checks) else 1


if __name__ == "__main__":
    sys.exit(main())

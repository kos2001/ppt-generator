#!/usr/bin/env python3
"""Programmatic grader for the Mode B (edit) evals — checks the objective
assertions in evals 3-5 against an output deck so results don't rest on
eyeballing a grid.

Reads only with python-pptx (no COM needed), so it grades a finished .pptx on
any platform even though producing it needed Windows PowerPoint.

Usage:
    python evals/grade_edit.py <output.pptx> --expect-slides 29
    python evals/grade_edit.py <output.pptx> --expect-slides 35 \
        --baseline decks/Claude_Cowork_No-Code_Automation.orig.pptx

Checks, reported as PASS/FAIL lines plus a JSON summary:
  * slide count equals --expect-slides (nothing deleted)
  * every slide carries a sequential 'i / N' page number
  * header-bar heights are consistent (a single dominant height -> the scaled
    chrome matches the native bars rather than stamping a fixed-size bar)
"""
import argparse
import json
import re
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PAGENO = re.compile(r"^\s*(\d+)\s*/\s*(\d+)\s*$")


def _cm(emu):
    return emu / 360000.0


def grade(path, expect_slides):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    n = len(prs.slides)
    results = []

    results.append((
        "slide count preserved (== %d, nothing deleted)" % expect_slides,
        n == expect_slides,
        "found %d slides" % n,
    ))

    nums, missing = [], []
    bar_heights = []
    for i, slide in enumerate(prs.slides, 1):
        page_num = None
        for shp in slide.shapes:
            if shp.has_text_frame:
                m = PAGENO.match(shp.text_frame.text.strip())
                if m:
                    page_num = (int(m.group(1)), int(m.group(2)))
            # a top-anchored full-width short rectangle == a header bar
            if (shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                    and (shp.top or 0) < 3000
                    and shp.width and shp.width > sw * 0.9
                    and shp.height and _cm(shp.height) < 5):
                bar_heights.append(round(_cm(shp.height), 2))
        if page_num is None:
            missing.append(i)
        else:
            nums.append((i, page_num))

    seq_ok = not missing and all(p == (i, n) for i, p in nums)
    results.append((
        "every slide numbered sequentially 'i / %d'" % n,
        seq_ok,
        "missing on %s" % missing if missing
        else "bad: %s" % [(i, p) for i, p in nums if p != (i, n)][:5] if not seq_ok
        else "all %d sequential" % n,
    ))

    uniq = sorted(set(bar_heights))
    # one dominant bar height (allow a tiny rounding spread) -> bars consistent
    consistent = len(uniq) <= 1 or (max(uniq) - min(uniq) <= 0.05)
    results.append((
        "header-bar heights consistent (scaled chrome matches native)",
        bool(bar_heights) and consistent,
        "bar heights cm: %s across %d bars" % (uniq, len(bar_heights)),
    ))

    return n, results


def main(argv=None):
    ap = argparse.ArgumentParser(description="Grade a Mode B edit output deck.")
    ap.add_argument("pptx")
    ap.add_argument("--expect-slides", type=int, required=True)
    args = ap.parse_args(argv)

    n, results = grade(args.pptx, args.expect_slides)
    passed = sum(1 for _, ok, _ in results if ok)
    for text, ok, evidence in results:
        print("[%s] %s — %s" % ("PASS" if ok else "FAIL", text, evidence))
    print(json.dumps({
        "pptx": args.pptx, "slides": n,
        "passed": passed, "total": len(results),
        "expectations": [
            {"text": t, "passed": ok, "evidence": e} for t, ok, e in results
        ],
    }, ensure_ascii=False))
    return 0 if passed == len(results) else 1


if __name__ == "__main__":
    sys.exit(main())
